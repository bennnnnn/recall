import io
from unittest.mock import MagicMock, patch

from app.services.attachment_content import (
    MAX_INDEX_EXTRACT_CHARS,
    bytes_match_claimed,
    extract_text_details,
)

XLSX = "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
PPTX = "application/vnd.openxmlformats-officedocument.presentationml.presentation"


def test_spreadsheet_values_coordinates_and_formulas():
    from openpyxl import Workbook

    book = Workbook()
    sheet = book.active
    sheet.title = "Budget"
    sheet.append(["Name", "Cost"])
    sheet.append(["Train", 42])
    sheet["D2"] = "=B2*2"
    output = io.BytesIO()
    book.save(output)
    data = output.getvalue()
    assert bytes_match_claimed(XLSX, data)
    assert not bytes_match_claimed(PPTX, data)
    details = extract_text_details(XLSX, data)
    assert details is not None
    assert "Budget" in details.text
    assert "B2: 42" in details.text
    assert "D2: =B2*2 (formula, not evaluated)" in details.text


def test_slide_text_table_and_notes():
    from pptx import Presentation
    from pptx.util import Inches

    deck = Presentation()
    slide = deck.slides.add_slide(deck.slide_layouts[1])
    slide.shapes.title.text = "Quarterly review"
    table = slide.shapes.add_table(1, 2, Inches(1), Inches(2), Inches(4), Inches(1)).table
    table.cell(0, 0).text = "Revenue"
    table.cell(0, 1).text = "42"
    slide.notes_slide.notes_text_frame.text = "Discuss retention"
    output = io.BytesIO()
    deck.save(output)
    data = output.getvalue()
    assert bytes_match_claimed(PPTX, data)
    details = extract_text_details(PPTX, data)
    assert details is not None
    assert "[slide 1]" in details.text
    assert "Quarterly review" in details.text
    assert "Revenue | 42" in details.text
    assert "Discuss retention" in details.text


def test_pdf_index_reaches_page_40_without_bloating_inline_excerpt():
    reader = MagicMock()
    reader.pages = [
        MagicMock(extract_text=MagicMock(return_value=f"Content {i}")) for i in range(1, 41)
    ]
    with patch("pypdf.PdfReader", return_value=reader):
        inline = extract_text_details("application/pdf", b"mock")
        indexed = extract_text_details(
            "application/pdf", b"mock", max_chars=MAX_INDEX_EXTRACT_CHARS
        )
    assert inline is not None and indexed is not None
    assert inline.page_capped
    assert "[page 40]" not in inline.text
    assert "[page 40] Content 40" in indexed.text
    assert not indexed.page_capped
