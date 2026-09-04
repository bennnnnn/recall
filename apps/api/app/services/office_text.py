"""Bounded, read-only Office excerpts. Never run formulas, macros or external links."""

from __future__ import annotations

import io
import logging
from collections.abc import Generator, Iterator
from typing import Any

logger = logging.getLogger(__name__)


def _spreadsheet_lines(data: bytes) -> Generator[str, None, None]:
    from openpyxl import load_workbook

    book = load_workbook(io.BytesIO(data), read_only=True, data_only=False, keep_links=False)
    try:
        yield "[Spreadsheet text: formulas are shown, not calculated; charts/images are not read.]"
        for sheet in book.worksheets:
            yield f"[sheet: {sheet.title}]"
            # Hard row/column bounds also cover falsely enormous worksheet dimensions.
            for row in sheet.iter_rows(
                max_row=min(sheet.max_row or 10000, 10000),
                max_col=min(sheet.max_column or 256, 256),
            ):
                values = []
                for cell in row:
                    if cell.value is not None:
                        suffix = " (formula, not evaluated)" if cell.data_type == "f" else ""
                        values.append(f"{cell.coordinate}: {cell.value}{suffix}")
                if values:
                    yield " | ".join(values)
            if (sheet.max_row or 0) > 10000 or (sheet.max_column or 0) > 256:
                yield (
                    "[File note: this sheet exceeds 10000 rows or 256 columns; "
                    "later cells were not read.]"
                )
    finally:
        book.close()


def _shape_text(shapes: Any, depth: int = 0) -> Iterator[str]:
    for shape in shapes:
        if shape.has_text_frame and shape.text.strip():
            yield str(shape.text)
        if shape.has_table:
            for row in shape.table.rows:
                yield " | ".join(cell.text for cell in row.cells)
        if depth < 8 and hasattr(shape, "shapes"):
            yield from _shape_text(shape.shapes, depth + 1)


def _presentation_lines(data: bytes) -> Generator[str, None, None]:
    from pptx import Presentation

    presentation = Presentation(io.BytesIO(data))
    yield "[Slide text and speaker notes only; embedded images/charts and animations are not read.]"
    for index, slide in enumerate(presentation.slides, 1):
        if index > 500:
            yield "[File note: slides after 500 were not read.]"
            break
        yield f"[slide {index}]"
        yield from _shape_text(slide.shapes)
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame is not None:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                yield f"[speaker notes] {notes}"


def extract_office_text(content_type: str, data: bytes, *, max_chars: int) -> str | None:
    lines = (
        _spreadsheet_lines(data)
        if content_type.endswith("spreadsheetml.sheet")
        else _presentation_lines(data)
    )
    parts: list[str] = []
    size = 0
    try:
        for line in lines:
            parts.append(line[: max_chars - size])
            size += len(parts[-1]) + 1
            if size >= max_chars:
                break
        return "\n".join(parts)[:max_chars] or None
    except Exception:
        logger.debug("Office text extraction failed", exc_info=True)
        return None
    finally:
        lines.close()
